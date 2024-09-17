from pptx import Presentation
import os
import tempfile

def generateSlides(source_path, new_pptx, file, verse):
    if file.startswith("~$"):
        # Skip temporary files
        return
    
    prs = Presentation(os.path.join(source_path, file))

    C_included = True

    # Handle the case where specific verses are selected
    if verse != "all":
        # Split the verse string into individual verse numbers
        verse_numbers = verse.split(',')
        for slide in prs.slides:
            # Extract the slide title to determine the verse number
            slide_title = slide.shapes.title.text
            # Split the slide title to extract the verse number
            parts = slide_title.split('-')
            if len(parts) < 2:
                continue  # Skip slides without a valid verse number
            slide_verse = parts[0]
            # Check if the verse number is in the list of selected verses
            if slide_verse in verse_numbers or (C_included and slide_verse == 'C'):
                C_included = True
                # Clone the slide and append it to the destination presentation
                new_slide = new_pptx.slides.add_slide(slide.slide_layout)
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        # Ignore placeholders (like titles and footers)
                        continue
                    if shape.shape_type == 13:  # Check if shape is a picture
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        image = shape.image
                        # Save the image to a temporary file
                        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                            temp_file.write(image.blob)
                            temp_file.close()
                            # Add the image to the slide from the temporary file
                            new_slide.shapes.add_picture(temp_file.name, left, top, width, height)
                            os.unlink(temp_file.name)
                    else:
                        # Check if the shape has a shape type
                        if shape.has_text_frame:
                            text = shape.text_frame.text
                        else:
                            text = None
                        # Create a new shape with the same properties as the original shape
                        if text:
                            new_shape = new_slide.shapes.add_textbox(
                                shape.left, shape.top, shape.width, shape.height
                            )
                            new_shape.text_frame.text = text
                        else:
                            new_shape = new_slide.shapes.add_shape(
                                shape.auto_shape_type, shape.left, shape.top,
                                shape.width, shape.height
                            )
                        # Copy the text from the source shape to the destination shape
                        if shape.has_text_frame:
                            new_shape.text_frame.text = shape.text_frame.text
            else:
                C_included = False
    # Handle the case where all slides are included
    else:
        for slide in prs.slides:
            # Clone the slide and append it to the destination presentation
            new_slide = new_pptx.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                if shape.is_placeholder:
                    # Ignore placeholders (like titles and footers)
                    continue
                if shape.shape_type == 13:  # Check if shape is a picture
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    image = shape.image
                    # Save the image to a temporary file
                    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                        temp_file.write(image.blob)
                        temp_file.close()
                        # Add the image to the slide from the temporary file
                        new_slide.shapes.add_picture(temp_file.name, left, top, width, height)
                        os.unlink(temp_file.name)
                else:
                    # Check if the shape has a shape type
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                    else:
                        text = None
                    # Create a new shape with the same properties as the original shape
                    if text:
                        new_shape = new_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        )
                        new_shape.text_frame.text = text
                    else:
                        new_shape = new_slide.shapes.add_shape(
                            shape.auto_shape_type, shape.left, shape.top,
                            shape.width, shape.height
                        )
                    # Copy the text from the source shape to the destination shape
                    if shape.has_text_frame:
                        new_shape.text_frame.text = shape.text_frame.text

    # Add a blank slide at the end
    new_pptx.slides.add_slide(new_pptx.slide_layouts[5])

def main():
    new_pptx = Presentation()

    source_path = input('Enter the slides path("/path/to/dir"): ')

    songs_w_verses = input('Enter songs and verses in the following format(nnnn-v,v,v nnnn-v,v nnnn-v,v)(nnnn-all if all verses/song does not have verses): ')

    arr = songs_w_verses.split()
    songs = []
    verses = []

    for item in arr:
        parts = item.split("-")
        songs.append(parts[0])
        verses.append(parts[1])

    for i in range(len(songs)):
        for _, _, files in os.walk(source_path):
            for file in files:
                if songs[i] in file:
                    generateSlides(source_path, new_pptx, file, verses[i])

    new_pptx.save("output.pptx")

if __name__ == "__main__":
    main()
